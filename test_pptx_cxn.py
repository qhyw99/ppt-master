from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Add two shapes
shape1 = slide.shapes.add_shape(1, Pt(100), Pt(100), Pt(100), Pt(100))
shape2 = slide.shapes.add_shape(1, Pt(300), Pt(100), Pt(100), Pt(100))

# Add a connector
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(200), Pt(150), Pt(300), Pt(150))

# Connect them
connector.begin_x = shape1.left + shape1.width
connector.begin_y = shape1.top + shape1.height / 2
connector.end_x = shape2.left
connector.end_y = shape2.top + shape2.height / 2

# python-pptx doesn't expose the glue point API natively easily without editing XML,
# but let's just see how a raw connector is built.
prs.save('/workspace/test_cxn.pptx')
