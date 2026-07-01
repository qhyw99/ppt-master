import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Add a group
# python-pptx doesn't have native add_group API, so we build it manually or just write raw XML
